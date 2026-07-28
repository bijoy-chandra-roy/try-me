"""Build TryMe 11-slide defense PPTX (shape-based diagrams, no remote render)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "defense" / "TryMe-Defense-Slides.pptx"

BG = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x5C, 0x5C)
ACCENT = RGBColor(0x5A, 0x6B, 0x3F)
RULE = RGBColor(0xE0, 0xDC, 0xD4)
PANEL = RGBColor(0xF7, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, text, size, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri", anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.paragraphs[0]  # ensure
        try:
            tf.auto_size = None
        except Exception:
            pass
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for run in p.runs:
        set_run(run, run.text, size, bold, color, font)
    return box


def add_bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        set_run(run, f"•  {item}", size, False, INK)
    return box


def add_notes(slide, notes):
    ns = slide.notes_slide.notes_text_frame
    ns.text = "\n".join(f"• {n}" for n in notes)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return slide


def speaker_tag(slide, name):
    add_textbox(slide, Inches(0.55), Inches(0.22), Inches(5), Inches(0.3), f"Speaker: {name}", size=11, color=MUTED)


def title(slide, text, top=Inches(0.48)):
    add_textbox(slide, Inches(0.55), top, Inches(12.2), Inches(0.5), text, size=26, bold=True)


def footer(slide):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.12), Inches(12.2), Pt(1.25))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(7.18), Inches(10), Inches(0.25), "TryMe — Sonargaon University Team Defense", size=10, color=MUTED)


def panel(slide, left, top, width, height, fill=PANEL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RULE
    shape.line.width = Pt(1)
    return shape


def pill(slide, left, top, width, height, text, fill=ACCENT, text_color=WHITE, size=12, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    try:
        shape.text_frame.paragraphs[0].space_before = Pt(0)
    except Exception:
        pass
    # vertical center approx via padding text
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            set_run(run, run.text, size, bold, text_color)
    return shape


def arrow_right(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.45), Inches(0.28))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def arrow_down(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.28), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ===== 1 Title =====
    s = blank_slide(prs)
    add_textbox(
        s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.1),
        "TryMe: Enterprise Virtual Try-On Architecture",
        size=34, bold=True, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, Inches(0.8), Inches(3.75), Inches(11.7), Inches(0.45),
        "Team Defense — Sonargaon University",
        size=18, color=MUTED, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s, Inches(0.8), Inches(4.45), Inches(11.7), Inches(0.4),
        "Part 1 · Product & Business Value  ·  Speaker: Sadman",
        size=13, color=ACCENT, align=PP_ALIGN.CENTER,
    )
    add_notes(s, [
        "Introduce the team and the product name",
        "TryMe closes the online-retail visualization gap with AI virtual try-on",
        "Three parts: product value → modeling → architecture/resilience",
    ])

    # ===== 2 Market =====
    s = blank_slide(prs)
    speaker_tag(s, "Sadman")
    title(s, "The Market Problem")
    panel(s, Inches(0.7), Inches(1.35), Inches(5.7), Inches(3.0))
    panel(s, Inches(6.9), Inches(1.35), Inches(5.7), Inches(3.0))
    add_textbox(s, Inches(0.7), Inches(2.3), Inches(5.7), Inches(0.55), "Uncertain Fit", size=24, bold=True, align=PP_ALIGN.CENTER, color=ACCENT)
    add_textbox(s, Inches(0.9), Inches(2.95), Inches(5.3), Inches(0.9), "Shoppers cannot map flat product photos to their own body", size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(6.9), Inches(2.3), Inches(5.7), Inches(0.55), "High Return Rates", size=24, bold=True, align=PP_ALIGN.CENTER, color=ACCENT)
    add_textbox(s, Inches(7.1), Inches(2.95), Inches(5.3), Inches(0.9), "Reverse logistics inflate cost and erode margin", size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.7), Inches(4.7), Inches(11.8), Inches(1.8), [
        "Visualization gap in online retail",
        "High reverse-logistics costs",
        "Conversion friction",
    ], size=18)
    footer(s)
    add_notes(s, [
        "Online shoppers cannot mentally map flat product photos to their body",
        "Returns create shipping, restocking, and markdown costs",
        "Friction at the confidence step kills conversion before checkout",
    ])

    # ===== 3 Solution flow =====
    s = blank_slide(prs)
    speaker_tag(s, "Sadman")
    title(s, "The TryMe Solution & User Flow")
    nodes = [
        (Inches(0.7), "User Reference\nPhoto"),
        (Inches(3.55), "Garment\nImage"),
        (Inches(6.55), "TryMe\nEngine"),
        (Inches(9.7), "Virtual\nPreview"),
    ]
    for left, label in nodes:
        pill(s, left, Inches(2.4), Inches(2.4), Inches(1.35), label, fill=PANEL, text_color=INK, size=14)
    arrow_right(s, Inches(3.15), Inches(2.9))
    # plus between photo and garment visually implied; engine receives both
    add_textbox(s, Inches(2.95), Inches(2.0), Inches(0.5), Inches(0.35), "+", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    arrow_right(s, Inches(6.05), Inches(2.9))
    arrow_right(s, Inches(9.1), Inches(2.9))
    add_bullets(s, Inches(0.7), Inches(4.5), Inches(11.8), Inches(1.8), [
        "Frictionless guest experience (Rate-limited: 3/hr)",
        "Seamless conversion to authenticated customer",
    ], size=18)
    footer(s)
    add_notes(s, [
        "Guest can try on without signing up, capped at 3/hour",
        "Authenticated customers get history, cart, checkout",
        "Output is a composite preview: person + selected garment",
    ])

    # ===== 4 SDLC =====
    s = blank_slide(prs)
    speaker_tag(s, "Yuvraj")
    title(s, "Software Development Life Cycle (SDLC)")
    add_bullets(s, Inches(0.55), Inches(1.15), Inches(12), Inches(1.1), [
        "Evolutionary Prototyping",
        "Risk-driven methodology",
        "Spirals: 1. Prototype, 2. Auth/RBAC, 3. Commerce, 4. Polish",
    ], size=15)
    phases = ["1. Objectives", "2. Risks", "3. Develop & Test", "4. Review & Plan"]
    for i, ph in enumerate(phases):
        left = Inches(0.7 + i * 3.1)
        pill(s, left, Inches(2.6), Inches(2.7), Inches(0.7), ph, size=12)
        if i < 3:
            arrow_right(s, left + Inches(2.75), Inches(2.8))
    add_textbox(s, Inches(0.7), Inches(3.5), Inches(12), Inches(0.35), "Each spiral cycle (then advances to the next)", size=12, color=MUTED)
    spirals = [
        "Spiral 1\nPrototype\nVTO + Circuit Breaker",
        "Spiral 2\nAuth / RBAC\n6 actors + dashboards",
        "Spiral 3\nCommerce\nCart, COD, orders",
        "Spiral 4\nPolish\nDesign, i18n, Vercel",
    ]
    for i, sp in enumerate(spirals):
        left = Inches(0.7 + i * 3.1)
        panel(s, left, Inches(4.0), Inches(2.85), Inches(2.4))
        add_textbox(s, left + Inches(0.1), Inches(4.25), Inches(2.65), Inches(2.0), sp, size=13, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_right(s, left + Inches(2.88), Inches(5.0))
    footer(s)
    add_notes(s, [
        "Spiral chosen because VTO API risk was unknown",
        "Each spiral ends with a working, deployable increment",
        "Spiral 1: VTO downtime → Circuit Breaker + fallback first",
    ])

    # ===== 5 Use cases =====
    s = blank_slide(prs)
    speaker_tag(s, "Yuvraj")
    title(s, "Behavioral Modeling")
    add_bullets(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.7), [
        "System boundaries defined",
        "6 Distinct Actors: Guest, Customer, Merchant, Support, Admin, Super Admin",
    ], size=14)
    actors = ["Guest", "Customer", "Merchant", "Support", "Admin", "Super Admin"]
    for i, a in enumerate(actors):
        pill(s, Inches(0.55 + i * 2.1), Inches(1.85), Inches(1.95), Inches(0.55), a, size=11)
    groups = [
        ("Catalog & Try-On", ["Browse catalog", "Filter category", "Virtual try-on", "Try-on history", "Live vs Fallback"]),
        ("Commerce", ["Manage cart", "Checkout (COD)", "Track orders", "Reviews", "Addresses"]),
        ("Account", ["Register / Sign in", "Google OAuth", "Profile & prefs"]),
        ("Merchant", ["Manage products", "Analytics", "Store profile"]),
        ("Admin", ["Users & roles", "Merchants", "Platform stats", "System config", "Assume role"]),
    ]
    for i, (g, items) in enumerate(groups):
        left = Inches(0.45 + i * 2.55)
        panel(s, left, Inches(2.65), Inches(2.4), Inches(4.0))
        add_textbox(s, left + Inches(0.08), Inches(2.75), Inches(2.25), Inches(0.4), g, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_bullets(s, left + Inches(0.1), Inches(3.2), Inches(2.2), Inches(3.2), items, size=11)
    footer(s)
    add_notes(s, [
        "Primary actors are six human roles; ImgBB, VTO, MongoDB, Google are secondary",
        "Guest is anonymous (browse + rate-limited try-on)",
        "Super Admin uniquely assumes roles via UC21",
    ])

    # ===== 6 Activity =====
    s = blank_slide(prs)
    speaker_tag(s, "Yuvraj")
    title(s, "Dynamic State & Workflow")
    add_bullets(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.65), [
        "Parallel control logic",
        "Asynchronous image storage and database querying",
    ], size=14)
    lanes = [
        ("User", ["Browse", "Try-on", "Cart", "Checkout", "Track"]),
        ("Frontend", ["Hooks/UI", "TryOnModal", "Badge", "Forms"]),
        ("Middleware", ["Auth?", "Permission?"]),
        ("Route Handlers", ["Products", "Try-on+RL", "Cart/COD", "Orders"]),
        ("Services", ["Product", "TryOn+CB", "Cart/Order"]),
        ("External", ["MongoDB", "ImgBB", "VTO SSE", "Fallback"]),
    ]
    for i, (lane, items) in enumerate(lanes):
        left = Inches(0.4 + i * 2.15)
        panel(s, left, Inches(1.85), Inches(2.05), Inches(4.7))
        add_textbox(s, left, Inches(1.95), Inches(2.05), Inches(0.35), lane, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_bullets(s, left + Inches(0.08), Inches(2.4), Inches(1.9), Inches(3.9), items, size=12)
    footer(s)
    add_notes(s, [
        "Swimlanes: User, Frontend, Middleware, API, Services, External",
        "Try-on fans out to ImgBB, VTO SSE, and optional MongoDB history",
        "Circuit breaker guarantees a preview when VTO fails",
    ])

    # ===== 7 ER =====
    s = blank_slide(prs)
    speaker_tag(s, "Yuvraj")
    title(s, "Data Architecture")
    add_bullets(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.65), [
        "MongoDB (Mongoose) physical collections",
        "Optimized relational references (User, Product, Order, TryOnHistory)",
    ], size=14)
    entities = [
        ("User", "email, role,\nmerchantId"),
        ("Merchant", "name, ownerId,\nstatus"),
        ("Product", "price, imageUrl,\nmerchantId"),
        ("Cart", "userId,\nitems[]"),
        ("Order", "userId,\nstatus, total"),
        ("Address", "userId,\nstreet"),
        ("Review", "user+product\n+order"),
        ("TryOnHistory", "user, product,\nfromFallback"),
        ("SystemConfig", "guestTryOnLimit\nmaintenance"),
    ]
    for i, (name, fields) in enumerate(entities):
        row, col = divmod(i, 5)
        left = Inches(0.55 + col * 2.5)
        top = Inches(1.9 + row * 2.35)
        panel(s, left, top, Inches(2.3), Inches(2.1))
        add_textbox(s, left, top + Inches(0.15), Inches(2.3), Inches(0.4), name, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.1), top + Inches(0.65), Inches(2.1), Inches(1.2), fields, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s)
    add_notes(s, [
        "Physical collections map 1:1 to Mongoose models",
        "TryOnHistory records fromFallback for resilience path",
        "SystemConfig singleton controls guest limit and maintenance",
    ])

    # ===== 8 Component =====
    s = blank_slide(prs)
    speaker_tag(s, "Bijoy")
    title(s, "Component Architecture")
    add_bullets(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.7), [
        "Unified Next.js 15 App Router (No standalone Express server)",
        "Feature-based directory structure (Route Handlers → Services → Repositories)",
    ], size=14)
    layers = [
        ("Client", "Pages · Features · Hooks\nEdge Middleware (JWT + RBAC)"),
        ("API", "Route Handlers\n/api/try-on · cart · orders · …"),
        ("Server", "Services → Repositories\nCircuit Breaker · VTO · ImgBB clients"),
        ("External", "MongoDB · ImgBB\nHF IDM-VTON · Google OAuth"),
    ]
    for i, (name, body) in enumerate(layers):
        top = Inches(1.95 + i * 1.15)
        panel(s, Inches(1.5), top, Inches(10.3), Inches(1.0))
        add_textbox(s, Inches(1.7), top + Inches(0.12), Inches(2.2), Inches(0.7), name, size=16, bold=True, color=ACCENT)
        add_textbox(s, Inches(4.0), top + Inches(0.15), Inches(7.5), Inches(0.7), body, size=14, color=INK)
        if i < 3:
            arrow_down(s, Inches(6.4), top + Inches(0.95))
    footer(s)
    add_notes(s, [
        "One Next.js app owns UI, middleware, and API route handlers",
        "Business logic in services; repositories isolate MongoDB",
        "External AI and image hosts behind clients + circuit breaker",
    ])

    # ===== 9 Circuit breaker =====
    s = blank_slide(prs)
    speaker_tag(s, "Bijoy")
    title(s, "Algorithmic Fault Tolerance")
    add_bullets(s, Inches(0.55), Inches(1.05), Inches(12), Inches(0.65), [
        "Mitigates third-party API network latency",
        "300-second timeout trips circuit to serve local cache fallback",
    ], size=14)
    # flowchart column
    steps = [
        (Inches(1.9), "Try-on request", ACCENT, WHITE),
        (Inches(2.7), "Upload image → ImgBB", PANEL, INK),
        (Inches(3.5), "Circuit Breaker execute", PANEL, INK),
        (Inches(4.3), "Race VTO SSE vs 300s timer", ACCENT, WHITE),
    ]
    for top, label, fill, tc in steps:
        pill(s, Inches(0.7), top, Inches(5.5), Inches(0.55), label, fill=fill, text_color=tc, size=12)
        if top != Inches(4.3):
            arrow_down(s, Inches(3.3), top + Inches(0.55))
    # branch
    panel(s, Inches(0.55), Inches(5.05), Inches(2.7), Inches(1.4))
    add_textbox(s, Inches(0.65), Inches(5.2), Inches(2.5), Inches(1.1), "Success\nfromFallback: false\nLive composite", size=12, align=PP_ALIGN.CENTER)
    panel(s, Inches(3.5), Inches(5.05), Inches(2.7), Inches(1.4))
    add_textbox(s, Inches(3.6), Inches(5.2), Inches(2.5), Inches(1.1), "Timeout / error\nfromFallback: true\nfallback-vto-result.jpg", size=12, align=PP_ALIGN.CENTER)
    # pseudocode
    panel(s, Inches(6.8), Inches(1.85), Inches(5.9), Inches(4.7))
    code = (
        "execute(operation):\n"
        "  try:\n"
        "    result = race(\n"
        "      operation(),\n"
        "      timeout(300_000 ms)\n"
        "    )\n"
        "    return { data: result,\n"
        "             fromFallback: false }\n"
        "  catch:\n"
        "    fallback = read(\n"
        "      \"cache/fallback-vto-result.jpg\")\n"
        "    return { data: fallback,\n"
        "             fromFallback: true }"
    )
    box = s.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(5.5), Inches(4.3))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        set_run(run, line, 12, False, INK, "Consolas")
    footer(s)
    add_notes(s, [
        "Free-tier Hugging Face Spaces are slow and can fail mid-SSE",
        "Timeout via VTO_TIMEOUT_MS; default 300000 ms (5 minutes)",
        "UI shows Live vs Fallback badge — demos always produce a result",
    ])

    # ===== 10 Security =====
    s = blank_slide(prs)
    speaker_tag(s, "Bijoy")
    title(s, "Critical Security Analysis")
    cols = [
        ("API Vulnerabilities", "Shared compute space risks (e.g., prompt / input injection on third-party inference hosts)"),
        ("Storage Instability", "ImgBB latency spikes and public-URL hosting dependency"),
        ("Data Compliance", "Third-party transmission of personal photos (user reference images leave the app boundary)"),
    ]
    for i, (h, body) in enumerate(cols):
        left = Inches(0.55 + i * 4.2)
        panel(s, left, Inches(1.5), Inches(3.95), Inches(4.8))
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.35), Inches(1.9), Inches(1.25), Inches(1.25))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        add_textbox(s, left + Inches(0.2), Inches(3.4), Inches(3.55), Inches(0.7), h, size=15, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.25), Inches(4.2), Inches(3.45), Inches(1.7), body, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s)
    add_notes(s, [
        "VTO runs on shared Hugging Face infrastructure",
        "ImgBB free host: weak availability and retention for PII imagery",
        "Production needs private storage, contracted AI hosts, consent policy",
    ])

    # ===== 11 Roadmap =====
    s = blank_slide(prs)
    speaker_tag(s, "Bijoy")
    title(s, "Future Roadmap")
    panel(s, Inches(0.55), Inches(1.2), Inches(5.7), Inches(3.3))
    panel(s, Inches(7.05), Inches(1.2), Inches(5.7), Inches(3.3))
    add_textbox(s, Inches(0.55), Inches(1.35), Inches(5.7), Inches(0.4), "Current Prototype", size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.05), Inches(1.35), Inches(5.7), Inches(0.4), "Production Evolution", size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.75), Inches(1.9), Inches(5.3), Inches(2.4), [
        "Next.js 15 App Router",
        "MongoDB + Mongoose",
        "ImgBB public URLs",
        "HF Space IDM-VTON (SSE)",
        "Circuit Breaker + fallback",
    ], size=13)
    add_bullets(s, Inches(7.25), Inches(1.9), Inches(5.3), Inches(2.4), [
        "Postgres + RLS (e.g. Supabase)",
        "Private buckets / signed URLs",
        "Self-hosted VTO on GPU",
        "Retain Circuit Breaker",
        "Stripe + CI/E2E + CDN + metrics",
    ], size=13)
    arrow_right(s, Inches(6.35), Inches(2.6))
    add_bullets(s, Inches(0.55), Inches(4.7), Inches(12.2), Inches(2.0), [
        "Infrastructure sovereignty for production retail VTO",
        "Supabase (Postgres + RLS + private buckets) replacing MongoDB / public-ImgBB coupling",
        "Self-hosted open-source VTO on dedicated GPU clusters replacing free Hugging Face Spaces",
        "Near-term productization: Stripe (or regional gateway), CI/E2E, CDN caching, VTO latency observability",
    ], size=12)
    footer(s)
    add_notes(s, [
        "Frame as productionization of a validated Spiral-4 prototype",
        "Layered architecture keeps storage and AI hosts swappable",
        "Circuit breaker remains valuable after self-hosting",
        "How: repo adapters + RLS; containerized VTON API; Stripe behind payment interface",
    ])

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
